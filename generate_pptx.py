"""
Génère la présentation PowerPoint du projet SentimentAI.
Usage : python generate_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette couleurs ───────────────────────────────────────────────────────────
DARK_BG     = RGBColor(0x0D, 0x11, 0x17)   # Fond sombre
ACCENT      = RGBColor(0x63, 0x66, 0xF1)   # Violet indigo
ACCENT2     = RGBColor(0x10, 0xB9, 0x81)   # Vert émeraude
ACCENT3     = RGBColor(0xF5, 0x9E, 0x0B)   # Orange ambre
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0x8B, 0x94, 0xA3)
CARD_BG     = RGBColor(0x16, 0x1B, 0x22)   # Fond carte


# ── Helpers ────────────────────────────────────────────────────────────────────

def set_slide_bg(slide, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_accent_bar(slide, color=None):
    """Barre décorative en haut de la slide."""
    c = color or ACCENT
    add_rect(slide, 0, 0, 10, 0.06, c)


def add_bullet_card(slide, title, bullets, left, top, width, height, accent=ACCENT):
    """Carte avec titre et liste de points."""
    # Fond carte
    add_rect(slide, left, top, width, height, CARD_BG, accent, Pt(1.5))
    # Titre carte
    add_text(slide, title, left + 0.15, top + 0.08, width - 0.3, 0.35,
             font_size=13, bold=True, color=accent)
    # Séparateur
    add_rect(slide, left + 0.15, top + 0.42, width - 0.3, 0.02, accent)
    # Bullets
    y = top + 0.52
    for b in bullets:
        add_text(slide, f"• {b}", left + 0.2, y, width - 0.4, 0.28,
                 font_size=10.5, color=WHITE)
        y += 0.28


def add_metric_box(slide, label, value, left, top, width=1.9, height=1.1, color=ACCENT):
    add_rect(slide, left, top, width, height, CARD_BG, color, Pt(2))
    add_text(slide, value, left, top + 0.08, width, 0.55,
             font_size=28, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, label, left, top + 0.63, width, 0.4,
             font_size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ── Slides ─────────────────────────────────────────────────────────────────────

def slide_titre(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Bande gauche colorée
    add_rect(slide, 0, 0, 0.12, 7.5, ACCENT)

    # Badge MLOps
    add_rect(slide, 1.0, 1.2, 2.2, 0.4, ACCENT)
    add_text(slide, "PROJET MLOPS END-TO-END", 1.0, 1.22, 2.2, 0.38,
             font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Titre principal
    add_text(slide, "SentimentAI", 1.0, 1.75, 8.5, 1.1,
             font_size=54, bold=True, color=WHITE)
    add_text(slide, "Analyse de sentiment · Trustpilot", 1.0, 2.8, 8.5, 0.6,
             font_size=22, color=ACCENT, bold=True)

    # Sous-titre
    add_text(slide,
             "API REST sécurisée · LightGBM + TF-IDF · Monitoring drift · MLflow",
             1.0, 3.5, 8.5, 0.5, font_size=14, color=LIGHT_GRAY)

    # Ligne séparatrice
    add_rect(slide, 1.0, 4.15, 7.5, 0.03, ACCENT)

    # Technos en bas
    techs = ["FastAPI", "LightGBM", "Streamlit", "Nginx", "Docker", "MLflow", "SHAP"]
    x = 1.0
    for t in techs:
        add_rect(slide, x, 4.35, 1.12, 0.35, CARD_BG, ACCENT, Pt(1))
        add_text(slide, t, x, 4.37, 1.12, 0.33,
                 font_size=9.5, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        x += 1.18

    # Métriques
    add_text(slide, "Accuracy : 71.95%   ·   F1-Score : 71.86%   ·   28 tests pytest   ·   565k avis analysés",
             1.0, 5.1, 8.5, 0.4, font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


def slide_contexte(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide)

    add_text(slide, "Contexte & Problématique", 0.4, 0.2, 9, 0.55,
             font_size=26, bold=True, color=WHITE)
    add_rect(slide, 0.4, 0.72, 1.2, 0.04, ACCENT)

    # Problème
    add_bullet_card(slide, "🎯  Le problème",
        ["Des milliers d'avis Trustpilot non structurés chaque jour",
         "Impossible de les traiter manuellement à grande échelle",
         "Besoin d'une classification automatique : Positif / Neutre / Négatif",
         "Et d'un système de production complet, pas juste un notebook"],
        0.3, 1.0, 4.4, 2.3, ACCENT)

    # Dataset
    add_bullet_card(slide, "📊  Le dataset",
        ["565 814 avis Trustpilot (Amazon)",
         "Notes 1-5 → regroupées en 3 classes",
         "277 826 lignes après nettoyage",
         "Dataset équilibré : 44 042 par classe",
         "Colonnes : reviewText, summary, overall"],
        4.9, 1.0, 4.4, 2.3, ACCENT2)

    # Labels
    boxes = [
        ("⭐ 1-2", "Négatif", ACCENT3),
        ("⭐ 3",   "Neutre",  LIGHT_GRAY),
        ("⭐ 4-5", "Positif", ACCENT2),
    ]
    x = 0.5
    for stars, label, color in boxes:
        add_rect(slide, x, 3.6, 2.8, 1.1, CARD_BG, color, Pt(2))
        add_text(slide, stars, x, 3.65, 2.8, 0.45,
                 font_size=20, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, label, x, 4.1, 2.8, 0.45,
                 font_size=14, color=WHITE, align=PP_ALIGN.CENTER)
        x += 3.0

    add_text(slide, "Objectif : pipeline MLOps complet, du preprocessing au monitoring en production",
             0.3, 5.0, 9.2, 0.45, font_size=12, color=LIGHT_GRAY, italic=True)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide)

    add_text(slide, "Architecture Technique", 0.4, 0.2, 9, 0.55,
             font_size=26, bold=True, color=WHITE)
    add_rect(slide, 0.4, 0.72, 1.5, 0.04, ACCENT)

    # Blocs architecture
    blocks = [
        (":8501", "Streamlit", "Interface utilisateur\nDark glassmorphism\n5 onglets", ACCENT, 0.3),
        (":8080", "Nginx", "Reverse proxy\nRate limit 2 req/s\nauth_request", ACCENT3, 2.15),
        (":80", "FastAPI", "Backend API\nAuth · Predict\nMonitoring", ACCENT2, 4.0),
        ("", "Models", "LightGBM .pkl\nTF-IDF .pkl\nRead-only", RGBColor(0xEC,0x48,0x99), 5.85),
        ("", "Data", "users.json\npredictions.jsonl\nVolumes Docker", LIGHT_GRAY, 7.7),
    ]

    for port, name, desc, color, x in blocks:
        add_rect(slide, x, 1.05, 1.7, 1.8, CARD_BG, color, Pt(2))
        if port:
            add_rect(slide, x, 1.05, 1.7, 0.3, color)
            add_text(slide, port, x, 1.06, 1.7, 0.28,
                     font_size=10, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
        add_text(slide, name, x, 1.38, 1.7, 0.4,
                 font_size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x, 1.78, 1.7, 0.95,
                 font_size=9.5, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # Flèches texte
    arrows = [(1.95, "→"), (3.8, "→"), (5.65, "→"), (7.5, "→")]
    for x, arrow in arrows:
        add_text(slide, arrow, x, 1.7, 0.3, 0.4,
                 font_size=20, color=ACCENT, align=PP_ALIGN.CENTER)

    # Flux prédiction
    add_rect(slide, 0.3, 3.1, 9.1, 0.04, ACCENT)
    add_text(slide, "Flux d'une prédiction", 0.3, 3.22, 9, 0.35,
             font_size=13, bold=True, color=ACCENT)

    steps = [
        "1. Client POST /predict\n+ X-API-Key",
        "2. Nginx : rate\nlimit check",
        "3. auth_request\n→ /verify_admin",
        "4. Token OK ?\nQuota OK ?",
        "5. TF-IDF → \nLightGBM",
        "6. Log JSONL\n+ quota++",
        "7. Réponse :\nsentiment + score",
    ]
    x = 0.2
    for i, s in enumerate(steps):
        color = ACCENT if i % 2 == 0 else CARD_BG
        border = ACCENT
        add_rect(slide, x, 3.65, 1.28, 0.9, color, border, Pt(1))
        add_text(slide, s, x + 0.05, 3.68, 1.18, 0.85,
                 font_size=8.5, color=WHITE, align=PP_ALIGN.CENTER)
        x += 1.33


def slide_ml(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide)

    add_text(slide, "Pipeline Machine Learning", 0.4, 0.2, 9, 0.55,
             font_size=26, bold=True, color=WHITE)
    add_rect(slide, 0.4, 0.72, 1.5, 0.04, ACCENT2)

    # Étapes pipeline
    steps = [
        ("📝", "Texte brut", "reviewText + summary"),
        ("🔡", "Lowercase\n+ nettoyage", "Ponctuation, chiffres"),
        ("✂️", "Tokenisation\nNLTK", "punkt tokenizer"),
        ("🚫", "Stopwords\nfiltrés", "English + custom"),
        ("📚", "Lemmatisation\nWordNet", "WordNetLemmatizer"),
        ("📊", "TF-IDF\n20 000 features", "ngram (1,2)"),
        ("🤖", "LightGBM\n3 classes", "early stopping 50"),
    ]
    x = 0.2
    for icon, title, sub in steps:
        add_rect(slide, x, 1.0, 1.3, 1.4, CARD_BG, ACCENT2, Pt(1.5))
        add_text(slide, icon, x, 1.05, 1.3, 0.4,
                 font_size=18, align=PP_ALIGN.CENTER, color=WHITE)
        add_text(slide, title, x, 1.45, 1.3, 0.5,
                 font_size=9, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
        add_text(slide, sub, x, 1.93, 1.3, 0.35,
                 font_size=7.5, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
        if x < 8.0:
            add_text(slide, "→", x + 1.28, 1.5, 0.25, 0.4,
                     font_size=14, color=ACCENT2, align=PP_ALIGN.CENTER)
        x += 1.38

    # Hyperparamètres
    add_bullet_card(slide, "⚙️  Hyperparamètres LightGBM",
        ["objective: multiclass (3 classes)",
         "num_leaves: 64 · max_depth: -1",
         "learning_rate: 0.05",
         "n_estimators: 1000 (early stopping: 50)",
         "colsample_bytree: 0.8 · subsample: 0.8"],
        0.3, 2.65, 4.4, 2.2, ACCENT2)

    # Résultats
    add_bullet_card(slide, "📈  Résultats",
        ["Accuracy globale : 71.95%",
         "F1-Score (weighted) : 71.86%",
         "Meilleur sur : Positif (F1 ≈ 0.74)",
         "Plus difficile : Neutre (F1 ≈ 0.68)",
         "Plafond naturel pour 3 classes proches"],
        4.9, 2.65, 4.4, 2.2, ACCENT3)

    # SHAP
    add_rect(slide, 0.3, 5.05, 9.1, 0.5, CARD_BG, ACCENT, Pt(1))
    add_text(slide, "🔍  SHAP TreeExplainer — Explicabilité : identifie les mots qui ont le plus influencé chaque prédiction",
             0.5, 5.1, 8.8, 0.4, font_size=11, color=WHITE)


def slide_securite(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide, ACCENT3)

    add_text(slide, "Sécurité & API", 0.4, 0.2, 9, 0.55,
             font_size=26, bold=True, color=WHITE)
    add_rect(slide, 0.4, 0.72, 1.0, 0.04, ACCENT3)

    add_bullet_card(slide, "🛡️  Nginx — Première ligne de défense",
        ["Rate limiting : 2 req/s par IP, burst de 5",
         "auth_request : /verify_admin avant chaque /predict",
         "Requête rejetée avant d'atteindre FastAPI si invalide",
         "Reverse proxy : cache la topologie interne"],
        0.3, 1.0, 4.4, 2.2, ACCENT3)

    add_bullet_card(slide, "🔐  Backend FastAPI — Authentification",
        ["Mots de passe : bcrypt (résistant brute force)",
         "Migration auto : anciens hash SHA-256 → bcrypt",
         "Tokens : secrets.token_hex(32) = 64 chars hex",
         "Expiration : 30 jours (token_expires_at)",
         "RBAC : role user (quota 5/j) vs admin (illimité)"],
        4.9, 1.0, 4.4, 2.2, ACCENT3)

    add_bullet_card(slide, "⚡  Race condition — threading.Lock",
        ["Problème : 2 requêtes simultanées pouvaient dépasser le quota",
         "Solution : verrou threading.Lock sur la mise à jour du quota",
         "Lecture → vérification → écriture atomique",
         "Protection côté API même si Nginx laisse passer"],
        0.3, 3.45, 4.4, 2.0, RGBColor(0xEC, 0x48, 0x99))

    add_bullet_card(slide, "✅  Validation Pydantic",
        ["username : 3-50 chars, alphanumeric + _ -",
         "password : 6-100 chars minimum",
         "role : Literal['user', 'admin'] uniquement",
         "text : 1-5000 chars, non vide (strip())",
         "feedback : Literal['correct', 'incorrect']"],
        4.9, 3.45, 4.4, 2.0, RGBColor(0xEC, 0x48, 0x99))

    add_text(slide, "🔒  CORS configuré · Logs structurés (logging) · Quota reset automatique à minuit",
             0.3, 5.65, 9.1, 0.4, font_size=11, color=LIGHT_GRAY, italic=True)


def slide_monitoring(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide, ACCENT2)

    add_text(slide, "Monitoring de Drift", 0.4, 0.2, 9, 0.55,
             font_size=26, bold=True, color=WHITE)
    add_rect(slide, 0.4, 0.72, 1.3, 0.04, ACCENT2)

    # KL divergence
    add_bullet_card(slide, "📡  Divergence KL (Kullback-Leibler)",
        ["Mesure l'écart entre distribution récente et distribution d'entraînement",
         "Distribution de référence : 1/3 par classe (Négatif / Neutre / Positif)",
         "Calculée sur les 7 derniers jours par défaut",
         "Stockée dans predictions_log.jsonl (append-only)"],
        0.3, 1.0, 5.8, 2.1, ACCENT2)

    # Seuils
    seuils = [
        ("🟢 Normal", "KL < 0.10", "Rien à faire", ACCENT2),
        ("🟡 Warning", "0.10 – 0.30", "Surveiller", ACCENT3),
        ("🔴 Critique", "KL ≥ 0.30", "Réentraîner !", RGBColor(0xEF, 0x44, 0x44)),
    ]
    x = 6.35
    for label, kl, action, color in seuils:
        add_rect(slide, x, 1.0, 3.0, 0.9, CARD_BG, color, Pt(2))
        add_text(slide, label, x + 0.1, 1.05, 2.8, 0.3,
                 font_size=12, bold=True, color=color)
        add_text(slide, f"{kl}  →  {action}", x + 0.1, 1.38, 2.8, 0.4,
                 font_size=10, color=LIGHT_GRAY)
        x = 6.35
        if label == "🟢 Normal":
            x = 6.35
        add_rect(slide, 6.35, 1.0 + seuils.index((label, kl, action, color)) * 0.95,
                 3.0, 0.9, CARD_BG, color, Pt(2))

    # Refaire proprement les seuils (éviter le bug de l'index dans la boucle)
    for shape in slide.shapes:
        pass

    y = 1.0
    for label, kl, action, color in seuils:
        add_rect(slide, 6.35, y, 3.0, 0.82, CARD_BG, color, Pt(2))
        add_text(slide, label, 6.5, y + 0.05, 2.8, 0.3,
                 font_size=11, bold=True, color=color)
        add_text(slide, f"{kl}  →  {action}", 6.5, y + 0.38, 2.8, 0.35,
                 font_size=9.5, color=LIGHT_GRAY)
        y += 0.88

    # Alerte confiance
    add_rect(slide, 0.3, 3.3, 9.1, 0.55, CARD_BG, ACCENT3, Pt(1.5))
    add_text(slide, "⚠️  Seuil de confiance : si confiance moyenne < 55% → alerte même en mode Warning",
             0.5, 3.38, 8.8, 0.38, font_size=11, color=ACCENT3)

    # Ce qui est loggé
    add_bullet_card(slide, "📝  Ce qui est enregistré à chaque prédiction",
        ["timestamp · username · text_preview (80 chars)",
         "prediction (Négatif/Neutre/Positif) · confidence",
         "class_id · feedback (correct/incorrect/null)",
         "Rotation automatique du fichier à 50 000 lignes (garde 80%)"],
        0.3, 4.05, 4.4, 1.85, ACCENT2)

    add_bullet_card(slide, "🎛️  Dashboard admin (onglet Monitor)",
        ["Statut drift en temps réel (🟢🟡🔴)",
         "Timeline journalière : volume + confiance",
         "Distribution récente vs distribution entraînement",
         "Accuracy basée sur les feedbacks utilisateurs"],
        4.9, 4.05, 4.4, 1.85, ACCENT2)


def slide_mlflow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide, RGBColor(0x00, 0xBC, 0xD4))

    MLFLOW_COLOR = RGBColor(0x00, 0xBC, 0xD4)

    add_text(slide, "MLflow — Experiment Tracking", 0.4, 0.2, 9, 0.55,
             font_size=26, bold=True, color=WHITE)
    add_rect(slide, 0.4, 0.72, 1.6, 0.04, MLFLOW_COLOR)

    # Workflow
    add_text(slide, "Workflow de réentraînement", 0.3, 0.92, 9, 0.35,
             font_size=13, bold=True, color=MLFLOW_COLOR)

    workflow = [
        ("📡 Monitor\nKL ≥ 0.30", "Drift\ndétecté"),
        ("🐍 train.py\n--data csv", "Entraîne\nLightGBM"),
        ("📊 MLflow\nlog_metrics()", "Logue tout\nautomatiquement"),
        ("🌐 UI :5001\nlocalhost", "Compare\nles runs"),
        ("📦 models/\n.pkl mis à jour", "Nouveau\nmodèle"),
        ("🔄 docker\nrestart api", "Déploiement\ninstantané"),
    ]
    x = 0.2
    for title, sub in workflow:
        add_rect(slide, x, 1.4, 1.52, 1.05, CARD_BG, MLFLOW_COLOR, Pt(1.5))
        add_text(slide, title, x, 1.45, 1.52, 0.6,
                 font_size=9, bold=True, color=MLFLOW_COLOR, align=PP_ALIGN.CENTER)
        add_text(slide, sub, x, 2.02, 1.52, 0.38,
                 font_size=8.5, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
        if x < 8.5:
            add_text(slide, "→", x + 1.5, 1.75, 0.25, 0.35,
                     font_size=16, color=MLFLOW_COLOR, align=PP_ALIGN.CENTER)
        x += 1.62

    # Ce que MLflow logue
    add_bullet_card(slide, "📋  Ce qui est loggé automatiquement",
        ["Paramètres : max_features, learning_rate, num_leaves, n_estimators…",
         "Métriques : accuracy, f1_weighted, f1 par classe, precision, best_iteration",
         "Artefacts : confusion_matrix.png, feature_importance.png, classification_report.txt",
         "Modèle versionné dans le registry : SentimentAI-LightGBM v1, v2, v3…"],
        0.3, 2.7, 5.5, 2.15, MLFLOW_COLOR)

    # Avantage clé
    add_bullet_card(slide, "💡  Pourquoi MLflow ?",
        ["Historique complet de chaque entraînement",
         "Rollback en 2 commandes si le nouveau modèle est moins bon",
         "Comparaison visuelle des runs dans l'UI",
         "Base SQLite locale — aucun service cloud requis",
         "Profil Docker séparé : ne tourne qu'à la demande"],
        5.95, 2.7, 3.7, 2.15, MLFLOW_COLOR)

    add_text(slide, "💻  Commandes : python training/train.py --data csv  ·  mlflow ui --backend-store-uri sqlite:///mlflow/mlflow.db",
             0.3, 5.08, 9.1, 0.42, font_size=9.5, color=LIGHT_GRAY, italic=True)


def slide_tests(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide, RGBColor(0xEC, 0x48, 0x99))

    PINK = RGBColor(0xEC, 0x48, 0x99)

    add_text(slide, "Tests & Qualité du Code", 0.4, 0.2, 9, 0.55,
             font_size=26, bold=True, color=WHITE)
    add_rect(slide, 0.4, 0.72, 1.3, 0.04, PINK)

    # Métrique principale
    add_metric_box(slide, "tests pytest", "28 / 28", 0.3, 1.0, 2.2, 1.1, ACCENT2)
    add_metric_box(slide, "taux de succès", "100%", 2.7, 1.0, 2.2, 1.1, ACCENT2)
    add_metric_box(slide, "modèle ML", "mocké", 5.1, 1.0, 2.2, 1.1, ACCENT3)
    add_metric_box(slide, "isolation", "tmp files", 7.5, 1.0, 2.0, 1.1, ACCENT)

    add_bullet_card(slide, "📁  test_auth.py — Authentification",
        ["Inscription : succès, doublon, username trop court, mdp trop court",
         "Validation rôle : Literal['user', 'admin'] — refus de 'superadmin'",
         "Validation username : caractères spéciaux refusés",
         "Connexion : succès, mauvais mdp, utilisateur inconnu, rôle admin",
         "/verify_admin : token valide, manquant, invalide"],
        0.3, 2.35, 4.4, 2.5, PINK)

    add_bullet_card(slide, "📁  test_predict.py — Prédiction & Monitoring",
        ["/predict : succès, sans token (401), token invalide (403)",
         "Texte vide ou blanc → 422 (Pydantic)",
         "Incrémentation quota après prédiction",
         "/quota/status : user (3/5) vs admin (unlimited)",
         "/predictions/history : vide puis après prédiction",
         "/feedback : introuvable (404), valeur invalide (422)",
         "/health : modèle chargé, version API"],
        4.9, 2.35, 4.4, 2.5, PINK)

    add_rect(slide, 0.3, 5.1, 9.1, 0.45, CARD_BG, ACCENT2, Pt(1))
    add_text(slide, "✅  Stratégie : fichiers tmp isolés par test · modèle LightGBM mocké (MagicMock) · patch.object sur les bons modules",
             0.5, 5.18, 8.8, 0.32, font_size=10, color=WHITE)


def slide_resultats(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide)

    add_text(slide, "Résultats & Métriques", 0.4, 0.2, 9, 0.55,
             font_size=26, bold=True, color=WHITE)
    add_rect(slide, 0.4, 0.72, 1.2, 0.04, ACCENT)

    # Métriques principales
    metrics = [
        ("Accuracy", "71.95%", ACCENT),
        ("F1-Score\nweighted", "71.86%", ACCENT2),
        ("Tests\npytest", "28 / 28", RGBColor(0xEC, 0x48, 0x99)),
        ("Quota\nutilisateurs", "5 / jour", ACCENT3),
        ("Token\nexpiration", "30 jours", LIGHT_GRAY),
    ]
    x = 0.2
    for label, val, color in metrics:
        add_metric_box(slide, label, val, x, 0.95, 1.85, 1.15, color)
        x += 1.92

    # Par classe
    add_text(slide, "Performance par classe", 0.3, 2.3, 9, 0.35,
             font_size=13, bold=True, color=ACCENT)

    classes = [
        ("😞 Négatif", "F1 : 0.70", "Precision : 0.69\nRecall : 0.71", RGBColor(0xEF, 0x44, 0x44)),
        ("😐 Neutre",  "F1 : 0.68", "Precision : 0.71\nRecall : 0.65", ACCENT3),
        ("😊 Positif", "F1 : 0.74", "Precision : 0.74\nRecall : 0.74", ACCENT2),
    ]
    x = 0.5
    for name, f1, details, color in classes:
        add_rect(slide, x, 2.75, 2.9, 1.4, CARD_BG, color, Pt(2))
        add_text(slide, name, x, 2.82, 2.9, 0.38,
                 font_size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, f1, x, 3.18, 2.9, 0.42,
                 font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, details, x, 3.6, 2.9, 0.45,
                 font_size=9.5, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
        x += 3.1

    # Stack technique résumé
    add_bullet_card(slide, "🏗️  Stack technique",
        ["Backend : FastAPI + Uvicorn + Pydantic v2",
         "Proxy : Nginx (rate limit + auth_request)",
         "ML : LightGBM + TF-IDF (20k features) + SHAP",
         "Frontend : Streamlit (dark glassmorphism, 5 onglets)",
         "Infra : Docker Compose (3 services) + MLflow (profil training)"],
        0.3, 4.3, 5.5, 2.1, ACCENT)

    add_bullet_card(slide, "🚀  Ce qui le distingue",
        ["Pas juste un modèle — une API de production complète",
         "Sécurité multicouche (Nginx + FastAPI + Pydantic)",
         "Monitoring drift automatique (KL divergence)",
         "Tests automatisés 100% passants",
         "MLflow pour tracer chaque réentraînement"],
        5.95, 4.3, 3.7, 2.1, ACCENT2)


def slide_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, 10, 0.06, ACCENT)
    add_rect(slide, 0, 7.44, 10, 0.06, ACCENT)

    add_text(slide, "SentimentAI", 0.5, 0.8, 9, 0.8,
             font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "Pipeline MLOps End-to-End", 0.5, 1.55, 9, 0.5,
             font_size=20, color=ACCENT, align=PP_ALIGN.CENTER, bold=True)

    add_rect(slide, 2, 2.2, 5.8, 0.03, ACCENT)

    pillars = [
        ("🤖", "Machine Learning", "LightGBM · TF-IDF · SHAP"),
        ("🛡️", "Sécurité", "bcrypt · Nginx · RBAC · CORS"),
        ("📡", "Monitoring", "Drift KL · Feedback · Logs"),
        ("🧪", "Qualité", "28 tests · Pydantic · Logging"),
        ("📊", "MLflow", "Versioning · Registry · UI"),
    ]
    x = 0.3
    for icon, title, sub in pillars:
        add_rect(slide, x, 2.4, 1.82, 1.55, CARD_BG, ACCENT, Pt(1.5))
        add_text(slide, icon, x, 2.45, 1.82, 0.48,
                 font_size=22, align=PP_ALIGN.CENTER, color=WHITE)
        add_text(slide, title, x, 2.9, 1.82, 0.38,
                 font_size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_text(slide, sub, x, 3.27, 1.82, 0.55,
                 font_size=8.5, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
        x += 1.89

    add_text(slide, "Perspectives", 0.5, 4.2, 9, 0.38,
             font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    persp = [
        "🗄️ PostgreSQL\nà la place de JSON",
        "🔔 Alertes drift\nautomatiques",
        "📦 Packaging\nDocker multi-arch",
        "☁️ Déploiement\nVPS / cloud",
        "🔄 Réentraînement\nautomatique CI/CD",
    ]
    x = 0.35
    for p in persp:
        add_rect(slide, x, 4.65, 1.8, 0.85, CARD_BG, LIGHT_GRAY, Pt(1))
        add_text(slide, p, x, 4.7, 1.8, 0.8,
                 font_size=9, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
        x += 1.87

    add_text(slide, "github.com/Lakdou/TrustPilot-webservice",
             0.5, 5.75, 9, 0.4,
             font_size=13, color=ACCENT, align=PP_ALIGN.CENTER, bold=True)
    add_text(slide, "Accuracy : 71.95%   ·   F1 : 71.86%   ·   28 tests ✅   ·   Docker   ·   MLflow",
             0.5, 6.25, 9, 0.38, font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ── Main ───────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    slides_fn = [
        slide_titre,
        slide_contexte,
        slide_architecture,
        slide_ml,
        slide_securite,
        slide_monitoring,
        slide_mlflow,
        slide_tests,
        slide_resultats,
        slide_conclusion,
    ]

    for fn in slides_fn:
        fn(prs)
        print(f"  ✅ Slide '{fn.__name__}' générée")

    output = "SentimentAI_Presentation.pptx"
    prs.save(output)
    print(f"\n🎉 Présentation sauvegardée : {output}")
    print(f"   {len(slides_fn)} slides · Format 10×7.5 pouces (16:9)")


if __name__ == "__main__":
    main()
